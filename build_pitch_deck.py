"""Build the Facility Trust Desk 3-minute pitch deck.

Run:
    python build_pitch_deck.py

Output:
    facility_trust_desk_pitch.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------- Design tokens ----------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x1B, 0x2E)      # background
NAVY_2 = RGBColor(0x13, 0x29, 0x44)    # card / panel
INK = RGBColor(0xF5, 0xF7, 0xFA)       # primary text
MUTED = RGBColor(0xA8, 0xB5, 0xC6)     # secondary text
AMBER = RGBColor(0xFF, 0xB4, 0x4D)     # accent
TEAL = RGBColor(0x4F, 0xD1, 0xC5)      # accent 2
RED = RGBColor(0xE2, 0x6B, 0x6B)
GREEN = RGBColor(0x6B, 0xCB, 0x77)
YELLOW = RGBColor(0xF4, 0xC4, 0x3C)
GREY = RGBColor(0x6E, 0x7C, 0x91)

FONT = "Calibri"
FONT_MONO = "Consolas"


# ---------- Helpers ----------------------------------------------------------


def add_blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.shadow.inherit = False
    return slide


def add_text(
    slide,
    text: str,
    *,
    left,
    top,
    width,
    height,
    size=18,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
    italic=False,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        text = [text]

    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_accent_bar(slide, top=Inches(0.55), color=AMBER):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(0.55), Inches(0.08)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_footer(slide, slide_num: int, total: int):
    add_text(
        slide,
        "Facility Trust Desk  ·  Track 1  ·  DAIS 2026",
        left=Inches(0.6),
        top=Inches(7.05),
        width=Inches(8),
        height=Inches(0.3),
        size=10,
        color=MUTED,
    )
    add_text(
        slide,
        f"{slide_num} / {total}",
        left=Inches(12.0),
        top=Inches(7.05),
        width=Inches(0.8),
        height=Inches(0.3),
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_eyebrow(slide, text: str, color=AMBER, top=Inches(0.7)):
    add_text(
        slide,
        text.upper(),
        left=Inches(0.6),
        top=top,
        width=Inches(10),
        height=Inches(0.35),
        size=12,
        bold=True,
        color=color,
    )


def add_title(slide, text: str, top=Inches(1.15), size=40):
    add_text(
        slide,
        text,
        left=Inches(0.6),
        top=top,
        width=Inches(12),
        height=Inches(1.2),
        size=size,
        bold=True,
        color=INK,
    )


def add_pill(slide, text: str, *, left, top, width, height, fill, ink=NAVY):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.5
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = ink
    return s


def add_card(slide, *, left, top, width, height, fill=NAVY_2, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.04
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    return s


# ---------- Slides -----------------------------------------------------------


def slide_title(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)

    # Big amber accent block on the left
    block = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H
    )
    block.line.fill.background()
    block.fill.solid()
    block.fill.fore_color.rgb = AMBER

    add_text(
        s,
        "TRACK 1  ·  DAIS 2026 HACKATHON",
        left=Inches(0.9),
        top=Inches(1.6),
        width=Inches(10),
        height=Inches(0.4),
        size=14,
        bold=True,
        color=AMBER,
    )
    add_text(
        s,
        "Facility Trust Desk",
        left=Inches(0.9),
        top=Inches(2.05),
        width=Inches(12),
        height=Inches(1.4),
        size=64,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        "Can this facility actually do what it claims?",
        left=Inches(0.9),
        top=Inches(3.4),
        width=Inches(12),
        height=Inches(0.8),
        size=28,
        color=MUTED,
        italic=True,
    )

    # tier strip at the bottom
    tiers = [
        ("STRONG", GREEN),
        ("PARTIAL", YELLOW),
        ("WEAK / SUSPICIOUS", AMBER),
        ("NO CLAIM", GREY),
    ]
    x = Inches(0.9)
    for label, color in tiers:
        add_pill(s, label, left=x, top=Inches(5.4), width=Inches(2.4), height=Inches(0.55), fill=color)
        x += Inches(2.6)

    add_text(
        s,
        "A trust signal for every facility-capability pair in the Indian healthcare dataset.",
        left=Inches(0.9),
        top=Inches(6.2),
        width=Inches(12),
        height=Inches(0.5),
        size=14,
        color=MUTED,
    )
    add_footer(s, n, total)


def slide_problem(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    add_accent_bar(s)
    add_eyebrow(s, "The problem")
    add_title(s, "Directories list claims. They don't grade them.")

    # Left: the scenario
    add_card(s, left=Inches(0.6), top=Inches(2.6), width=Inches(6.0), height=Inches(3.9))
    add_text(
        s,
        "A district planner.",
        left=Inches(0.95),
        top=Inches(2.8),
        width=Inches(5.4),
        height=Inches(0.45),
        size=18,
        bold=True,
        color=AMBER,
    )
    add_text(
        s,
        [
            "A high-risk delivery is being referred.",
            "Minutes to choose a hospital.",
            "Searches \"maternity\".",
            "Gets 200 hospitals, all claiming it.",
        ],
        left=Inches(0.95),
        top=Inches(3.3),
        width=Inches(5.4),
        height=Inches(3.0),
        size=20,
        color=INK,
    )

    # Right: the gap, big quote
    add_card(s, left=Inches(7.0), top=Inches(2.6), width=Inches(5.7), height=Inches(3.9), fill=NAVY_2)
    add_text(
        s,
        "The question that matters:",
        left=Inches(7.3),
        top=Inches(2.85),
        width=Inches(5.2),
        height=Inches(0.5),
        size=14,
        color=MUTED,
    )
    add_text(
        s,
        "\u201cCan they actually do it?\u201d",
        left=Inches(7.3),
        top=Inches(3.4),
        width=Inches(5.2),
        height=Inches(1.5),
        size=36,
        bold=True,
        color=AMBER,
    )
    add_text(
        s,
        "No one tells the planner. Self-reported capability is treated as ground truth across the entire sector.",
        left=Inches(7.3),
        top=Inches(5.0),
        width=Inches(5.2),
        height=Inches(1.4),
        size=16,
        color=INK,
    )
    add_footer(s, n, total)


def slide_what(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    add_accent_bar(s)
    add_eyebrow(s, "What we built")
    add_title(s, "A trust signal for every facility \u00d7 capability pair.")

    add_text(
        s,
        "~10,000 Indian healthcare facilities  \u00d7  6 capabilities  \u2192  ~60,000 graded decisions",
        left=Inches(0.6),
        top=Inches(2.5),
        width=Inches(12.1),
        height=Inches(0.6),
        size=20,
        color=MUTED,
    )

    # Capability pills
    caps = ["ICU", "NICU", "Maternity", "Emergency", "Oncology", "Trauma"]
    x = Inches(0.6)
    for c in caps:
        add_pill(s, c, left=x, top=Inches(3.25), width=Inches(1.95), height=Inches(0.55), fill=TEAL)
        x += Inches(2.05)

    # Four tier cards
    tiers = [
        ("STRONG", GREEN, "Structured evidence and well-corroborated by independent sources."),
        ("PARTIAL", YELLOW, "One real signal. Not enough to fully confirm the claim."),
        ("WEAK / SUSPICIOUS", AMBER, "Prose-only, screening-only, implausible scale, or contradictory."),
        ("NO CLAIM", GREY, "Nothing in the facility's own evidence supports this capability."),
    ]
    x = Inches(0.6)
    card_w = Inches(3.0)
    for label, color, desc in tiers:
        card = add_card(s, left=x, top=Inches(4.4), width=card_w, height=Inches(2.3))
        # color stripe
        stripe = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, Inches(4.4), card_w, Inches(0.18)
        )
        stripe.line.fill.background()
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        add_text(
            s,
            label,
            left=x + Inches(0.2),
            top=Inches(4.75),
            width=card_w - Inches(0.4),
            height=Inches(0.45),
            size=16,
            bold=True,
            color=color,
        )
        add_text(
            s,
            desc,
            left=x + Inches(0.2),
            top=Inches(5.25),
            width=card_w - Inches(0.4),
            height=Inches(1.4),
            size=13,
            color=INK,
        )
        x += card_w + Inches(0.13)
    add_footer(s, n, total)


def slide_user_flow(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    add_accent_bar(s)
    add_eyebrow(s, "What the planner sees")
    add_title(s, "Rank. Drill into evidence. Override with a note.")

    steps = [
        ("1", "PICK", "Capability and state."),
        ("2", "RANK", "A ranked list of facilities with a visible tier."),
        ("3", "DRILL", "Why we made the call: structured matches, prose hits, sources, snippets."),
        ("4", "OVERRIDE", "Local note + reviewer tag, kept in the browser."),
    ]
    x = Inches(0.6)
    w = Inches(3.0)
    for num, head, body in steps:
        add_card(s, left=x, top=Inches(2.6), width=w, height=Inches(3.6))
        # big number
        add_text(
            s,
            num,
            left=x + Inches(0.2),
            top=Inches(2.7),
            width=Inches(1.2),
            height=Inches(1.2),
            size=64,
            bold=True,
            color=AMBER,
        )
        add_text(
            s,
            head,
            left=x + Inches(0.2),
            top=Inches(4.05),
            width=w - Inches(0.4),
            height=Inches(0.5),
            size=18,
            bold=True,
            color=TEAL,
        )
        add_text(
            s,
            body,
            left=x + Inches(0.2),
            top=Inches(4.6),
            width=w - Inches(0.4),
            height=Inches(1.5),
            size=14,
            color=INK,
        )
        x += w + Inches(0.13)

    add_text(
        s,
        "Planners can disagree with us. They cannot be lied to about what we know.",
        left=Inches(0.6),
        top=Inches(6.45),
        width=Inches(12.1),
        height=Inches(0.4),
        size=15,
        italic=True,
        color=AMBER,
    )
    add_footer(s, n, total)


def slide_pipeline(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    add_accent_bar(s)
    add_eyebrow(s, "Under the hood")
    add_title(s, "A deterministic pipeline with an auditable LLM second opinion.")

    # Pipeline boxes
    nodes = [
        ("RAW", "10,088 raw\nfacility rows", GREY),
        ("SILVER", "Cleaned, deduped,\npincode-repaired", TEAL),
        ("GOLD\n(heuristic)", "Tier from structured\n+ claim + prose + corroboration", AMBER),
        ("LLM\nSUB-SIGNALS", "Typed, fingerprinted,\nnever the final tier", AMBER),
        ("GOLD\n(final mart)", "Tier deterministically\nrecomputed in SQL", GREEN),
    ]
    n_nodes = len(nodes)
    total_w = Inches(12.1)
    gap = Inches(0.25)
    node_w = (total_w - gap * (n_nodes - 1)) / n_nodes
    node_h = Inches(2.0)
    top = Inches(2.85)
    x = Inches(0.6)
    centers = []
    for label, body, color in nodes:
        card = add_card(s, left=x, top=top, width=node_w, height=node_h)
        # accent top
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, node_w, Inches(0.14))
        stripe.line.fill.background()
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        add_text(
            s,
            label,
            left=x + Inches(0.1),
            top=top + Inches(0.3),
            width=node_w - Inches(0.2),
            height=Inches(0.7),
            size=15,
            bold=True,
            color=color,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            s,
            body,
            left=x + Inches(0.1),
            top=top + Inches(1.05),
            width=node_w - Inches(0.2),
            height=Inches(0.9),
            size=11,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
        centers.append(x + node_w / 2)
        x += node_w + gap

    # arrows between boxes
    arrow_y = top + node_h + Inches(0.15)
    for i in range(n_nodes - 1):
        # use a thin chevron between the cards
        ax = centers[i] + node_w / 2 - Inches(0.05)
        aw = gap + Inches(0.1)
        arr = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            ax - aw / 2 + node_w / 2,
            top + node_h / 2 - Inches(0.15),
            aw,
            Inches(0.3),
        )
        arr.line.fill.background()
        arr.fill.solid()
        arr.fill.fore_color.rgb = MUTED

    # Three principles below
    principles = [
        ("Deterministic tier", "Final tier is recomputed by SQL every run. Repeatable, auditable.", GREEN),
        ("LLM as second opinion", "Writes typed sub-signals only. Tier rules read them, then decide.", TEAL),
        ("Fingerprinted reviews", "Each LLM review is keyed to a hash of the evidence it saw.", AMBER),
    ]
    x = Inches(0.6)
    pw = (Inches(12.1) - Inches(0.4)) / 3
    py = Inches(5.55)
    for head, body, color in principles:
        add_card(s, left=x, top=py, width=pw, height=Inches(1.2))
        add_text(
            s,
            head,
            left=x + Inches(0.2),
            top=py + Inches(0.1),
            width=pw - Inches(0.4),
            height=Inches(0.4),
            size=14,
            bold=True,
            color=color,
        )
        add_text(
            s,
            body,
            left=x + Inches(0.2),
            top=py + Inches(0.5),
            width=pw - Inches(0.4),
            height=Inches(0.7),
            size=11,
            color=INK,
        )
        x += pw + Inches(0.2)
    add_footer(s, n, total)


def slide_discipline(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    add_accent_bar(s, color=RED)
    add_eyebrow(s, "What we deliberately did not ship", color=RED)
    add_title(s, "We killed our citation-quality score.")

    # Two columns: before / after
    add_card(s, left=Inches(0.6), top=Inches(2.6), width=Inches(6.0), height=Inches(3.9), fill=NAVY_2, border=RED)
    add_text(
        s,
        "TEMPTING",
        left=Inches(0.85),
        top=Inches(2.8),
        width=Inches(5.5),
        height=Inches(0.4),
        size=13,
        bold=True,
        color=RED,
    )
    add_text(
        s,
        "citation_support_quality: 0.82",
        left=Inches(0.85),
        top=Inches(3.3),
        width=Inches(5.5),
        height=Inches(0.6),
        size=22,
        bold=True,
        font=FONT_MONO,
        color=INK,
    )
    add_text(
        s,
        [
            "A confident-looking number per claim.",
            "Backed by\u2026 a list of homepage URLs.",
            "No per-claim snippets to verify it.",
            "A score with no source is theater.",
        ],
        left=Inches(0.85),
        top=Inches(4.0),
        width=Inches(5.5),
        height=Inches(2.4),
        size=15,
        color=INK,
    )

    add_card(s, left=Inches(7.0), top=Inches(2.6), width=Inches(5.7), height=Inches(3.9), fill=NAVY_2, border=GREEN)
    add_text(
        s,
        "WHAT WE SHIP",
        left=Inches(7.25),
        top=Inches(2.8),
        width=Inches(5.2),
        height=Inches(0.4),
        size=13,
        bold=True,
        color=GREEN,
    )
    add_text(
        s,
        "supporting_snippets[]",
        left=Inches(7.25),
        top=Inches(3.3),
        width=Inches(5.2),
        height=Inches(0.6),
        size=22,
        bold=True,
        font=FONT_MONO,
        color=INK,
    )
    add_text(
        s,
        [
            "Actual quotes from the facility's own text.",
            "Citation count and source mix as features.",
            "\u201cno_claim\u201d is a first-class tier.",
            "We'd rather show no number than a fake one.",
        ],
        left=Inches(7.25),
        top=Inches(4.0),
        width=Inches(5.2),
        height=Inches(2.4),
        size=15,
        color=INK,
    )
    add_footer(s, n, total)


def slide_by_the_numbers(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    add_accent_bar(s)
    add_eyebrow(s, "By the numbers")
    add_title(s, "Cleaned, joined, scored, served.")

    stats = [
        ("10,088", "raw facility rows"),
        ("9,932", "unique facilities after dedupe"),
        ("~60k", "facility \u00d7 capability decisions"),
        ("6", "capabilities graded"),
        ("4", "trust tiers"),
        ("10", "SQL steps, raw \u2192 silver \u2192 gold"),
    ]
    cols = 3
    cell_w = (Inches(12.1) - Inches(0.4)) / cols
    cell_h = Inches(1.8)
    for i, (num, label) in enumerate(stats):
        row, col = divmod(i, cols)
        x = Inches(0.6) + col * (cell_w + Inches(0.2))
        y = Inches(2.7) + row * (cell_h + Inches(0.2))
        add_card(s, left=x, top=y, width=cell_w, height=cell_h)
        add_text(
            s,
            num,
            left=x + Inches(0.2),
            top=y + Inches(0.15),
            width=cell_w - Inches(0.4),
            height=Inches(1.0),
            size=44,
            bold=True,
            color=AMBER,
        )
        add_text(
            s,
            label,
            left=x + Inches(0.2),
            top=y + Inches(1.15),
            width=cell_w - Inches(0.4),
            height=Inches(0.5),
            size=14,
            color=MUTED,
        )

    add_text(
        s,
        "Deployed as a Databricks App with the Analytics plugin against a live SQL Warehouse.",
        left=Inches(0.6),
        top=Inches(6.55),
        width=Inches(12.1),
        height=Inches(0.4),
        size=14,
        italic=True,
        color=MUTED,
    )
    add_footer(s, n, total)


def slide_ambition(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    add_accent_bar(s, color=TEAL)
    add_eyebrow(s, "Where this goes", color=TEAL)
    add_title(s, "Anywhere people act on claims someone made about themselves.")

    domains = [
        ("Government services", "Eligibility & service-level claims"),
        ("Suppliers & vendors", "Certifications, capacity, compliance"),
        ("Schools", "Programs, accreditations, outcomes"),
        ("NGOs", "Reach, impact, capability"),
    ]
    x = Inches(0.6)
    w = Inches(3.0)
    for head, body in domains:
        add_card(s, left=x, top=Inches(2.7), width=w, height=Inches(2.0))
        add_text(
            s,
            head,
            left=x + Inches(0.2),
            top=Inches(2.85),
            width=w - Inches(0.4),
            height=Inches(0.55),
            size=18,
            bold=True,
            color=TEAL,
        )
        add_text(
            s,
            body,
            left=x + Inches(0.2),
            top=Inches(3.5),
            width=w - Inches(0.4),
            height=Inches(1.1),
            size=13,
            color=INK,
        )
        x += w + Inches(0.13)

    # The pattern
    add_card(s, left=Inches(0.6), top=Inches(5.0), width=Inches(12.1), height=Inches(1.8), fill=NAVY_2, border=AMBER)
    add_text(
        s,
        "THE PATTERN",
        left=Inches(0.85),
        top=Inches(5.15),
        width=Inches(11.5),
        height=Inches(0.35),
        size=12,
        bold=True,
        color=AMBER,
    )
    add_text(
        s,
        "Deterministic ranking + an auditable LLM second opinion + an honest \u201cwe don\u2019t know\u201d state.",
        left=Inches(0.85),
        top=Inches(5.5),
        width=Inches(11.5),
        height=Inches(1.2),
        size=22,
        bold=True,
        color=INK,
    )
    add_footer(s, n, total)


def _appendix_eyebrow(s, label: str):
    add_accent_bar(s, color=TEAL)
    add_eyebrow(s, f"Appendix \u00b7 {label}", color=TEAL)


def slide_architecture(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    _appendix_eyebrow(s, "Architecture")
    add_title(s, "What runs where.")

    # Four-lane horizontal architecture: Sources -> Heuristic SQL -> LLM sidecar -> Serving
    lane_top = Inches(2.4)
    lane_h = Inches(4.0)
    lane_gap = Inches(0.18)
    total_w = Inches(12.1)
    lane_w = (total_w - lane_gap * 3) / 4
    xs = [Inches(0.6) + i * (lane_w + lane_gap) for i in range(4)]

    lane_headers = [
        ("SOURCES", "Unity Catalog (read-only)", MUTED),
        ("HEURISTIC SQL", "Databricks SQL Warehouse", AMBER),
        ("LLM SIDECAR", "Python runner + Model Serving", TEAL),
        ("SERVING", "Databricks App + Browser", GREEN),
    ]

    for x, (head, sub, color) in zip(xs, lane_headers):
        add_card(s, left=x, top=lane_top, width=lane_w, height=lane_h)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, lane_top, lane_w, Inches(0.14))
        stripe.line.fill.background()
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        add_text(
            s, head,
            left=x + Inches(0.15), top=lane_top + Inches(0.25),
            width=lane_w - Inches(0.3), height=Inches(0.4),
            size=13, bold=True, color=color, align=PP_ALIGN.CENTER,
        )
        add_text(
            s, sub,
            left=x + Inches(0.15), top=lane_top + Inches(0.7),
            width=lane_w - Inches(0.3), height=Inches(0.35),
            size=10, color=MUTED, align=PP_ALIGN.CENTER,
        )

    def lane_items(x, items, *, start_top=lane_top + Inches(1.2)):
        item_h = Inches(0.42)
        gap = Inches(0.1)
        y = start_top
        for label, sub, color in items:
            box = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x + Inches(0.12), y, lane_w - Inches(0.24), item_h,
            )
            box.adjustments[0] = 0.25
            box.line.color.rgb = color
            box.line.width = Pt(0.75)
            box.fill.solid()
            box.fill.fore_color.rgb = NAVY
            add_text(
                s, label,
                left=x + Inches(0.22), top=y + Inches(0.02),
                width=lane_w - Inches(0.44), height=Inches(0.22),
                size=10, bold=True, color=color, font=FONT_MONO,
            )
            add_text(
                s, sub,
                left=x + Inches(0.22), top=y + Inches(0.22),
                width=lane_w - Inches(0.44), height=Inches(0.2),
                size=8.5, color=MUTED,
            )
            y += item_h + gap

    lane_items(xs[0], [
        ("facilities", "10,088 raw rows", MUTED),
        ("india_post_pincode_directory", "pincode \u2192 district", MUTED),
        ("nfhs_5_district_health_indicators", "district health stats", MUTED),
    ])
    lane_items(xs[1], [
        ("01-03  silver.*_clean", "clean + dedupe", AMBER),
        ("04-05  gold.geo / need_index", "joins + planner context", AMBER),
        ("06  gold.*_assessment_heuristic", "deterministic baseline tier", AMBER),
    ])
    lane_items(xs[2], [
        ("07  silver.llm_inputs", "prompt-ready + fingerprint", TEAL),
        ("run_llm_capability_review.py", "Databricks Serving or OpenAI", TEAL),
        ("08  silver.llm_outputs_raw", "append-only model output", TEAL),
        ("09  gold.llm_signals", "typed sub-signals + status", TEAL),
    ])
    lane_items(xs[3], [
        ("10  gold.*_assessment", "final tier, SQL-recomputed", GREEN),
        ("Databricks App (AppKit)", "states / capability_ranked / facility_detail", GREEN),
        ("Browser (React + localStorage)", "ranking, drill-down, overrides", GREEN),
    ])

    # Arrows between lanes
    for i in range(3):
        arr = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            xs[i] + lane_w - Inches(0.05),
            lane_top + lane_h / 2 - Inches(0.15),
            lane_gap + Inches(0.1),
            Inches(0.3),
        )
        arr.line.fill.background()
        arr.fill.solid()
        arr.fill.fore_color.rgb = MUTED

    # Bottom guarantee strip
    add_card(s, left=Inches(0.6), top=Inches(6.55), width=Inches(12.1), height=Inches(0.4), fill=NAVY_2, border=AMBER)
    add_text(
        s,
        "The model never writes the final tier. Script 10 rebuilds it deterministically from script 06's baseline + current LLM sub-signals.",
        left=Inches(0.85), top=Inches(6.6), width=Inches(11.6), height=Inches(0.3),
        size=11, italic=True, color=INK, align=PP_ALIGN.CENTER,
    )
    add_footer(s, n, total)


def slide_worked_example(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    _appendix_eyebrow(s, "Worked example")
    add_title(s, "Greenfield General: oncology row, before and after LLM.")

    add_text(
        s,
        "Specialties list \"Oncology\". 6 sources, official site, named staff, 80 beds. "
        "Description says \"Our oncology screening clinic runs every Monday.\"",
        left=Inches(0.6), top=Inches(2.25), width=Inches(12.1), height=Inches(0.6),
        size=14, italic=True, color=MUTED,
    )

    col_top = Inches(2.95)
    col_h = Inches(3.95)
    col_w = (Inches(12.1) - Inches(0.3)) / 2

    # LEFT: Heuristic only
    x = Inches(0.6)
    add_card(s, left=x, top=col_top, width=col_w, height=col_h, border=AMBER)
    add_text(
        s, "HEURISTIC ONLY (script 06)",
        left=x + Inches(0.25), top=col_top + Inches(0.18),
        width=col_w - Inches(0.5), height=Inches(0.4),
        size=13, bold=True, color=AMBER,
    )
    add_text(
        s,
        [
            "structured_hit = true   (Oncology in specialties)",
            "claim_hit      = true",
            "prose_hit      = true",
            "screening_only = false  (regex saw \"oncology\")",
            "well_corroborated = true",
        ],
        left=x + Inches(0.25), top=col_top + Inches(0.65),
        width=col_w - Inches(0.5), height=Inches(1.7),
        size=12, color=INK, font=FONT_MONO,
    )
    add_text(
        s, "Tier rule that fires:",
        left=x + Inches(0.25), top=col_top + Inches(2.3),
        width=col_w - Inches(0.5), height=Inches(0.3),
        size=11, color=MUTED,
    )
    add_text(
        s, "structured + well_corroborated  \u2192  STRONG",
        left=x + Inches(0.25), top=col_top + Inches(2.6),
        width=col_w - Inches(0.5), height=Inches(0.45),
        size=14, bold=True, color=INK, font=FONT_MONO,
    )
    add_pill(s, "STRONG  ·  score 96",
             left=x + Inches(0.25), top=col_top + Inches(3.2),
             width=Inches(3.5), height=Inches(0.55), fill=GREEN)

    # RIGHT: Heuristic + LLM
    x2 = Inches(0.6) + col_w + Inches(0.3)
    add_card(s, left=x2, top=col_top, width=col_w, height=col_h, border=GREEN)
    add_text(
        s, "HEURISTIC + LLM (script 10)",
        left=x2 + Inches(0.25), top=col_top + Inches(0.18),
        width=col_w - Inches(0.5), height=Inches(0.4),
        size=13, bold=True, color=GREEN,
    )
    add_text(
        s,
        [
            "LLM read the prose. Returned typed signals:",
            "  screening_only_llm   = true",
            "  capability_scope_llm = screening_or_diagnostics_only",
            "  llm_review_status    = current",
            "structured_hit, implausible, well_corroborated unchanged.",
        ],
        left=x2 + Inches(0.25), top=col_top + Inches(0.65),
        width=col_w - Inches(0.5), height=Inches(1.7),
        size=11, color=INK, font=FONT_MONO,
    )
    add_text(
        s, "Tier rule that fires (NEW, top of ladder):",
        left=x2 + Inches(0.25), top=col_top + Inches(2.3),
        width=col_w - Inches(0.5), height=Inches(0.3),
        size=11, color=MUTED,
    )
    add_text(
        s, "capability_scope = screening_only  \u2192  WEAK",
        left=x2 + Inches(0.25), top=col_top + Inches(2.6),
        width=col_w - Inches(0.5), height=Inches(0.45),
        size=14, bold=True, color=INK, font=FONT_MONO,
    )
    add_pill(s, "WEAK / SUSPICIOUS  ·  score 36",
             left=x2 + Inches(0.25), top=col_top + Inches(3.2),
             width=Inches(4.2), height=Inches(0.55), fill=AMBER)
    add_footer(s, n, total)


def slide_two_ladders(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    _appendix_eyebrow(s, "The tier ladders")
    add_title(s, "Script 10 is script 06's ladder with two new rules on top.")

    col_top = Inches(2.55)
    col_h = Inches(4.2)
    col_w = (Inches(12.1) - Inches(0.3)) / 2

    rules_06 = [
        ("1", "no evidence", "no_claim", MUTED),
        ("2", "oncology + screening_only + NOT structured", "weak", AMBER),
        ("3", "implausible", "weak", AMBER),
        ("4", "structured + well_corroborated", "strong", GREEN),
        ("5", "structured OR claim", "partial", YELLOW),
        ("6", "prose only", "weak", AMBER),
        ("7", "ELSE", "no_claim", MUTED),
    ]
    rules_10 = [
        ("1", "no evidence", "no_claim", MUTED, False),
        ("2", "capability_scope = screening_or_diagnostics_only", "weak", TEAL, True),
        ("3", "capability_scope = adjacent_service", "weak", TEAL, True),
        ("4", "oncology + screening_only + NOT structured", "weak", AMBER, False),
        ("5", "implausible", "weak", AMBER, False),
        ("6", "structured + well_corroborated", "strong", GREEN, False),
        ("7", "structured OR claim", "partial", YELLOW, False),
        ("8", "prose only", "weak", AMBER, False),
        ("9", "ELSE", "no_claim", MUTED, False),
    ]

    def draw_ladder(x, header, rules, header_color):
        add_card(s, left=x, top=col_top, width=col_w, height=col_h)
        add_text(
            s, header,
            left=x + Inches(0.25), top=col_top + Inches(0.15),
            width=col_w - Inches(0.5), height=Inches(0.4),
            size=14, bold=True, color=header_color,
        )
        row_h = Inches(0.4)
        gap = Inches(0.06)
        y = col_top + Inches(0.7)
        for row in rules:
            if len(row) == 5:
                num, condition, result, color, is_new = row
            else:
                num, condition, result, color = row
                is_new = False
            # background highlight for NEW rules
            if is_new:
                hl = s.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x + Inches(0.18), y - Inches(0.02),
                    col_w - Inches(0.36), row_h + Inches(0.04),
                )
                hl.adjustments[0] = 0.2
                hl.line.color.rgb = TEAL
                hl.line.width = Pt(0.75)
                hl.fill.solid()
                hl.fill.fore_color.rgb = NAVY
            add_text(
                s, num,
                left=x + Inches(0.25), top=y + Inches(0.04),
                width=Inches(0.3), height=row_h,
                size=12, bold=True, color=MUTED, font=FONT_MONO,
            )
            add_text(
                s, condition,
                left=x + Inches(0.6), top=y + Inches(0.04),
                width=col_w - Inches(2.3), height=row_h,
                size=11, color=INK, font=FONT_MONO,
            )
            add_text(
                s, "\u2192 " + result,
                left=x + col_w - Inches(1.65), top=y + Inches(0.04),
                width=Inches(1.4), height=row_h,
                size=11, bold=True, color=color, font=FONT_MONO, align=PP_ALIGN.RIGHT,
            )
            if is_new:
                add_text(
                    s, "NEW",
                    left=x + col_w - Inches(0.75), top=y + Inches(0.04),
                    width=Inches(0.5), height=row_h,
                    size=9, bold=True, color=TEAL,
                )
            y += row_h + gap

    draw_ladder(Inches(0.6), "SCRIPT 06  ·  heuristic only", rules_06, AMBER)
    draw_ladder(Inches(0.6) + col_w + Inches(0.3), "SCRIPT 10  ·  heuristic + LLM", rules_10, GREEN)

    add_text(
        s,
        "First matching rule wins. The two NEW rules only fire when the LLM review is current; "
        "otherwise capability_scope is NULL and rules 2 and 3 cannot match, so script 10 behaves like script 06.",
        left=Inches(0.6), top=Inches(6.85), width=Inches(12.1), height=Inches(0.4),
        size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER,
    )
    add_footer(s, n, total)


def slide_why_llm_overrides(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)
    _appendix_eyebrow(s, "Why the LLM can do what the regex can't")
    add_title(s, "Same idea. Two safety guards removed, on purpose.")

    # Three guard cards
    card_top = Inches(2.5)
    card_h = Inches(3.4)
    card_w = (Inches(12.1) - Inches(0.4)) / 3
    xs = [Inches(0.6) + i * (card_w + Inches(0.2)) for i in range(3)]

    cards = [
        (
            "GUARD 1  (06)", AMBER,
            "screening_only is a regex",
            [
                "Flags when text has",
                "  cancer / tumor",
                "AND NOT",
                "  oncolog%, chemotherap%,",
                "  cancer treatment, ...",
                "",
                "Greenfield text contains",
                "\"oncology\", so the flag is",
                "false. Regex cannot tell",
                "\"oncology screening clinic\"",
                "from \"oncology department\".",
            ],
        ),
        (
            "GUARD 2  (06)", AMBER,
            "Refuses to overrule structured fields",
            [
                "Rule requires",
                "  NOT structured_hit",
                "",
                "If \"Oncology\" is in the formal",
                "specialty list, the heuristic",
                "trusts that fact over its own",
                "fuzzy regex.",
                "",
                "Greenfield has structured_hit",
                "= true, so the rule cannot fire",
                "even if the regex had triggered.",
            ],
        ),
        (
            "LLM RULE  (10)", GREEN,
            "Both guards removed",
            [
                "capability_scope =",
                "  screening_or_diagnostics_only",
                "",
                "Typed verdict from the model",
                "after reading the actual prose.",
                "Explicitly allowed to overrule",
                "the structured specialty list.",
                "",
                "Safer because the LLM read the",
                "page, not a regex.",
            ],
        ),
    ]

    for x, (label, color, head, body) in zip(xs, cards):
        add_card(s, left=x, top=card_top, width=card_w, height=card_h, border=color)
        add_text(
            s, label,
            left=x + Inches(0.2), top=card_top + Inches(0.15),
            width=card_w - Inches(0.4), height=Inches(0.35),
            size=11, bold=True, color=color,
        )
        add_text(
            s, head,
            left=x + Inches(0.2), top=card_top + Inches(0.5),
            width=card_w - Inches(0.4), height=Inches(0.6),
            size=15, bold=True, color=INK,
        )
        add_text(
            s, body,
            left=x + Inches(0.2), top=card_top + Inches(1.15),
            width=card_w - Inches(0.4), height=card_h - Inches(1.3),
            size=11, color=INK, font=FONT_MONO,
        )

    # Safety strip: fingerprint + fallback
    add_card(s, left=Inches(0.6), top=Inches(6.05), width=Inches(12.1), height=Inches(0.95), fill=NAVY_2, border=TEAL)
    add_text(
        s, "SAFETY NET",
        left=Inches(0.85), top=Inches(6.15), width=Inches(11.5), height=Inches(0.3),
        size=11, bold=True, color=TEAL,
    )
    add_text(
        s,
        "Every LLM review is tied to a hash of the evidence it saw. If the facility text changes the "
        "review goes \"stale\"; script 10 then falls back to script 06's heuristic flags automatically. "
        "If no LLM review exists at all, script 10 produces the same answer as script 06.",
        left=Inches(0.85), top=Inches(6.45), width=Inches(11.5), height=Inches(0.55),
        size=11, color=INK,
    )
    add_footer(s, n, total)


def slide_thanks(prs: Presentation, n: int, total: int):
    s = add_blank_slide(prs)

    block = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H
    )
    block.line.fill.background()
    block.fill.solid()
    block.fill.fore_color.rgb = AMBER

    add_text(
        s,
        "THANK YOU",
        left=Inches(0.9),
        top=Inches(2.4),
        width=Inches(12),
        height=Inches(0.5),
        size=18,
        bold=True,
        color=AMBER,
    )
    add_text(
        s,
        "Facility Trust Desk",
        left=Inches(0.9),
        top=Inches(2.9),
        width=Inches(12),
        height=Inches(1.2),
        size=56,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        "Questions?",
        left=Inches(0.9),
        top=Inches(4.15),
        width=Inches(12),
        height=Inches(0.8),
        size=28,
        italic=True,
        color=MUTED,
    )
    add_text(
        s,
        "Raw  \u2192  Silver  \u2192  Gold  \u2192  LLM sub-signals  \u2192  Deterministic final tier",
        left=Inches(0.9),
        top=Inches(5.4),
        width=Inches(12),
        height=Inches(0.5),
        size=14,
        color=MUTED,
        font=FONT_MONO,
    )
    add_text(
        s,
        "Databricks App  \u00b7  SQL Warehouse  \u00b7  Llama on Databricks Model Serving",
        left=Inches(0.9),
        top=Inches(5.85),
        width=Inches(12),
        height=Inches(0.4),
        size=13,
        color=MUTED,
    )
    add_footer(s, n, total)


# ---------- Build ------------------------------------------------------------


def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_what,
        slide_user_flow,
        slide_pipeline,
        slide_discipline,
        slide_by_the_numbers,
        slide_ambition,
        slide_thanks,
        slide_architecture,
        slide_worked_example,
        slide_two_ladders,
        slide_why_llm_overrides,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        fn(prs, i, total)

    prs.save(out_path)
    print(f"Wrote {out_path} ({total} slides).")


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "facility_trust_desk_pitch.pptx"
    build(out)
